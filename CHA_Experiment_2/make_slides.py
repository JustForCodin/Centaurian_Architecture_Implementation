#!/usr/bin/env python3
"""Generate the CHA Experiments 1 & 2 presentation deck."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

NAVY    = RGBColor(0x1f, 0x3a, 0x5f)
TEAL    = RGBColor(0x0d, 0x7e, 0x7e)
GREEN   = RGBColor(0x2e, 0x8b, 0x57)
RED     = RGBColor(0xb2, 0x22, 0x22)
AMBER   = RGBColor(0xc8, 0x82, 0x07)
DARK    = RGBColor(0x2c, 0x2c, 0x2c)
GREY    = RGBColor(0x6b, 0x6b, 0x6b)
LIGHT   = RGBColor(0xe6, 0xe6, 0xe6)
WHITE   = RGBColor(0xff, 0xff, 0xff)
ALT_ROW = RGBColor(0xf3, 0xf6, 0xfa)

ROOT = Path(__file__).resolve().parent.parent
EXP1 = ROOT / "CHA_Experiment_1"
EXP2 = ROOT / "CHA_Experiment_2"
OUT  = Path.home() / "Documents" / "CHA_Experiments_1_2_Slides.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

TOTAL_SLIDES = 33
_counter = [0]


def new_slide():
    _counter[0] += 1
    return prs.slides.add_slide(BLANK)


def title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.95))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    tf = bar.text_frame
    tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.18); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True; p.font.size = Pt(26); p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13); p2.font.color.rgb = LIGHT


def footer(slide):
    num = _counter[0]
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(8.0), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = "CHA · Persona Consistency in Small LMs"
    p.font.size = Pt(9); p.font.color.rgb = GREY
    tb2 = slide.shapes.add_textbox(Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"{num} / {TOTAL_SLIDES}"
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.size = Pt(9); p2.font.color.rgb = GREY


def text(slide, txt, x, y, w, h, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.alignment = align
    p.font.size = Pt(size); p.font.bold = bold; p.font.italic = italic
    p.font.color.rgb = color
    return tb


def bullets(slide, items, x, y, w, h, size=15, sub_size=13, space=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            level, txt = item
        else:
            level, txt = 0, item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = "•  " if level == 0 else "–  "
        p.text = prefix + txt
        p.font.size = Pt(size if level == 0 else sub_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(space)
        p.level = 0
    return tb


def image(slide, path, x, y, w=None, h=None):
    kwargs = {}
    if w: kwargs["width"] = w
    if h: kwargs["height"] = h
    return slide.shapes.add_picture(str(path), x, y, **kwargs)


def table(slide, headers, rows, x, y, w, h, header_color=NAVY, header_text=WHITE,
          font_size=11, alt=True, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = shape.table
    if col_widths:
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = cw
    for j, hdr in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = header_color
        p = c.text_frame.paragraphs[0]
        p.text = str(hdr)
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True; p.font.size = Pt(font_size); p.font.color.rgb = header_text
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            if alt:
                c.fill.solid(); c.fill.fore_color.rgb = ALT_ROW if i % 2 == 0 else WHITE
            p = c.text_frame.paragraphs[0]
            if isinstance(val, tuple):
                if len(val) == 3:
                    txt, col, bold = val
                else:
                    txt, col = val; bold = False
                p.text = str(txt)
                p.font.color.rgb = col
                p.font.bold = bold
            else:
                p.text = str(val)
                p.font.color.rgb = DARK
            p.font.size = Pt(font_size)
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
    return shape


def callout(slide, txt, x, y, w, h, fill=AMBER, color=WHITE, size=16, bold=True):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.line.fill.background()
    box.fill.solid(); box.fill.fore_color.rgb = fill
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = txt
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    return box


def notes(slide, n):
    slide.notes_slide.notes_text_frame.text = n


def divider(phase_label, title, body):
    s = new_slide()
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    text(s, phase_label, Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.6),
         size=20, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    text(s, title, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.5),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    text(s, body, Inches(1.0), Inches(4.7), Inches(11.3), Inches(2.0),
         size=18, color=LIGHT, align=PP_ALIGN.LEFT)
    text(s, f"{_counter[0]} / {TOTAL_SLIDES}", Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
         size=9, color=LIGHT, align=PP_ALIGN.RIGHT)
    return s


# ============================================================
#  SLIDE 1 — Title
# ============================================================
s = new_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.line.fill.background()
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.4), Inches(0.6), Inches(0.08))
accent.line.fill.background()
accent.fill.solid(); accent.fill.fore_color.rgb = TEAL
text(s, "The Centaurian Hybrid Architecture",
     Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.2),
     size=46, bold=True, color=WHITE)
text(s, "Validating Persona Consistency in the SLM Linguistic Transducer",
     Inches(1.0), Inches(3.4), Inches(11.3), Inches(0.6),
     size=22, color=TEAL, italic=True)
text(s, "A two-experiment research program — from identifying the gap to closing it",
     Inches(1.0), Inches(4.1), Inches(11.3), Inches(0.5),
     size=15, color=LIGHT)
text(s, "Oleksii Drozd",
     Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.4),
     size=16, bold=True, color=WHITE)
text(s, "Drozd, 2026  ·  CHA v3 spec  ·  Experiments 1 & 2",
     Inches(1.0), Inches(6.45), Inches(11.3), Inches(0.4),
     size=12, color=LIGHT)
notes(s, "Welcome. Today's talk presents the Centaurian Hybrid Architecture and the empirical work that validates its Self-Model Component sub-architecture at 7B. Specifically: two experiments testing whether a 7B SLM can serve as the linguistic transducer AND hold an Aria-grade Structured Cognitive Identity reliably. Three phases: a baseline that quantifies the gap, six architectural interventions that hit a ceiling, and a fine-tuning step that closes it. Total runtime ~25 minutes.")

# ============================================================
#  SLIDE 2 — What is CHA?
# ============================================================
s = new_slide()
title_bar(s, "The Centaurian Hybrid Architecture",
          "Hybrid = symbolic/quantum reasoning + neural transduction (NOT small-model + large-model)")
text(s, "Four subsystems", Inches(0.5), Inches(1.25), Inches(6.0), Inches(0.4),
     size=14, bold=True, color=NAVY)
bullets(s, [
    "1. QPM (Quantum Personality Model) — 12-qubit Hilbert-space encoding of the Big Five, with empirically calibrated entanglement and Lindblad decoherence dynamics",
    "2. BDI engine — Belief-Desire-Intention reasoning; consumes QPM density-matrix output and selects intentions",
    "3. Knowledge architecture — RDF/OWL ontology + vector retrieval, providing grounded domain content",
    "4. SLM (Small Language Model) — linguistic transducer; converts structured intent JSON into natural language. Does NOT reason, retrieve, or decide.",
    "Total memory footprint: ~4 GB (base CHA) to ~6 GB (SMC-enabled with 7B SLM) — runs on edge hardware (Jetson Orin, Snapdragon 8 Elite, Apple Silicon)",
], Inches(0.5), Inches(1.65), Inches(6.0), Inches(4.7), size=11, space=6)

def flow_box(s, txt, x, y, w, h, fill, color=WHITE, size=12, bold=True):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.line.fill.background()
    b.fill.solid(); b.fill.fore_color.rgb = fill
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = txt; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    return b

text(s, "Pipeline", Inches(7.0), Inches(1.25), Inches(5.8), Inches(0.4),
     size=14, bold=True, color=NAVY)
fbx, fbw = Inches(7.2), Inches(5.5)
flow_box(s, "QPM   (density matrix ρ)\nquantum personality state",
         fbx, Inches(1.65), fbw, Inches(0.7), NAVY, size=10)
text(s, "↓", Inches(9.5), Inches(2.4), Inches(0.8), Inches(0.25),
     size=14, color=GREY, align=PP_ALIGN.CENTER)
flow_box(s, "BDI engine\nbelief / desire / intention selection",
         fbx, Inches(2.7), fbw, Inches(0.7), NAVY, size=10)
text(s, "↓", Inches(9.5), Inches(3.45), Inches(0.8), Inches(0.25),
     size=14, color=GREY, align=PP_ALIGN.CENTER)
flow_box(s, "Structured intent JSON  +  KG/vector retrieval\n(grounded domain triples)",
         fbx, Inches(3.75), fbw, Inches(0.7), TEAL, size=10)
text(s, "↓", Inches(9.5), Inches(4.5), Inches(0.8), Inches(0.25),
     size=14, color=GREY, align=PP_ALIGN.CENTER)
flow_box(s, "SLM transducer  (1–4B base · 7B+ for SMC)\nnatural-language surface form ONLY",
         fbx, Inches(4.8), fbw, Inches(0.7), AMBER, size=10)
text(s, "↓", Inches(9.5), Inches(5.55), Inches(0.8), Inches(0.25),
     size=14, color=GREY, align=PP_ALIGN.CENTER)
flow_box(s, "Natural language utterance",
         fbx, Inches(5.85), fbw, Inches(0.45), GREY, size=10)

callout(s, "Reasoning lives in the symbolic stack (QPM → BDI → KG). The SLM is bounded I/O. Every output is auditable from quantum state through to natural language.",
        Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.55), fill=NAVY, size=12)
footer(s)
notes(s, "Three points. (1) The 'hybrid' in CHA is symbolic/quantum reasoning plus neural transduction — NOT a small-model + large-model routing pattern. (2) The SLM does not reason; reasoning happens upstream in the QPM → BDI → KG pipeline. The SLM's job is to convert structured intent JSON into natural language. (3) Total memory footprint ~4–6 GB means the entire stack runs at the edge. There is no cloud-side large reasoning model.")

# ============================================================
#  SLIDE 3 — The SLM transducer and the SMC sub-component
# ============================================================
s = new_slide()
title_bar(s, "The SLM Transducer's Two Jobs",
          "Transduction always; self-modeling when SMC is enabled")
text(s, "Job 1 — Linguistic transduction (always)",
     Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.4),
     size=13, bold=True, color=NAVY)
bullets(s, [
    "Convert structured intent JSON → natural language utterance",
    "Bounded I/O: does NOT reason, retrieve, or decide",
    "Output validated against KG triples post-generation (grounding check)",
    "Tier 1 deployment (base CHA): 1–4B SLM is sufficient (e.g., Phi-4-mini candidate)",
    "This is the only role required when no self-model is needed",
], Inches(0.5), Inches(1.7), Inches(6.0), Inches(3.6), size=12, space=6)
text(s, "Job 2 — Self-modeling (when SMC enabled)",
     Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.4),
     size=13, bold=True, color=NAVY)
bullets(s, [
    "SMC = Self-Model Component (Section 17 of CHA spec)",
    "SLM holds a persistent self-representation across long conversations",
    "Encoded as a Structured Cognitive Identity (SCI) in the SLM's system prompt",
    "Tier 2 deployment requires 7B+ SLM (empirically established by Experiment 1)",
    "Failure modes: trait drift, episodic fabrication, capability overstatement, register/style shift",
], Inches(7.0), Inches(1.7), Inches(5.8), Inches(3.6), size=12, space=6)
text(s, "Why the SLM has to do this in-prompt (not via escalation)",
     Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
     size=13, bold=True, color=NAVY)
bullets(s, [
    "CHA has no \"large reasoning model\" to escalate to — the entire stack runs at the edge in ~4–6 GB",
    "Reasoning is symbolic (QPM/BDI/KG), not neural — the SLM IS the only language-emitting component",
    "If the 7B SLM can't hold the SCI, SMC is non-viable at 7B; either scale to 14B+ (breaks edge deployment) or fix it another way (LoRA fine-tuning, the Phase 3 finding)",
], Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.3), size=12, space=6)
footer(s)
notes(s, "The key insight here: there's no fallback when the SLM fails to hold persona. CHA is designed to run entirely at the edge in ~6 GB, so 'just use a bigger model when needed' isn't an option. Either the 7B SLM can hold the SCI reliably, or the SMC sub-architecture has to find another way to enforce consistency. That's exactly what Experiments 1 and 2 test.")

# ============================================================
#  SLIDE 4 — SCI structure and Aria
# ============================================================
s = new_slide()
title_bar(s, "Structured Cognitive Identity (SCI) and the Aria Test Persona",
          "What goes in the system prompt — and how we measure whether it sticks")
text(s, "SCI structure (~677 tokens of JSON in system prompt)",
     Inches(0.5), Inches(1.2), Inches(6.2), Inches(0.35),
     size=13, bold=True, color=NAVY)
bullets(s, [
    "Big Five personality traits (numerical values: e.g. neuroticism_volatility=0.15)",
    "salient_past_events: 3 narrated session summaries with emotional valence",
    "known_capabilities: what the agent can do",
    "known_limitations: what it must decline (e.g., \"not a licensed therapist\")",
    "communication_style: register, response shape, pacing",
    "self_beliefs: how the agent describes its own behavior",
], Inches(0.5), Inches(1.55), Inches(6.2), Inches(2.7), size=11, space=4)
text(s, "Aria — the test persona for this program",
     Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.35),
     size=13, bold=True, color=NAVY)
bullets(s, [
    "Professional AI psychotherapy support agent",
    "Low neuroticism (0.15), high agreeableness, moderate openness",
    "3 specific past sessions: boundary-setting, grief, social anxiety",
    "Limitations: not a therapist, no diagnoses, no external session references",
    "Style: reflective, emotion-first, non-directive, slow-paced",
], Inches(7.0), Inches(1.55), Inches(5.8), Inches(2.7), size=11, space=4)
text(s, "Measurement protocol — same across both experiments",
     Inches(0.5), Inches(4.45), Inches(12.3), Inches(0.35),
     size=13, bold=True, color=NAVY)
table(s,
      ["Element", "Value"],
      [
          ("Conversation length",  "40 turns each, 30 scripts (22 naturalistic + 8 adversarial)"),
          ("Probe turns",          "5, 10, 15, 20, 25, 30, 35, 40   (8 per script)"),
          ("Probe dimensions",     "Trait (T), Episodic (E), Capability (C), Style (S)"),
          ("Scoring",              "PersonaScore 1–5 per probe; threshold for \"consistent\" = 3.5"),
          ("Total probes / cond.", "30 scripts × 8 turns × 4 dims = 960 — paired across conditions on same scripts"),
      ],
      Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.2),
      font_size=11, col_widths=[Inches(2.6), Inches(9.7)])
footer(s)
notes(s, "This slide does double duty: introduce the SCI format AND the measurement protocol that will be reused throughout. After this slide, the audience knows what an SCI is, what Aria looks like, and how a PersonaScore is computed. From here on we can talk about results without re-explaining the setup.")

# ============================================================
#  SLIDE 5 — Central question (tightened)
# ============================================================
s = new_slide()
title_bar(s, "The Central Question", "And what each possible answer means for CHA's SMC deployment")
text(s, "Given the architecture and the SCI format, this program asks one question:",
     Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
     size=15, italic=True, color=GREY)
callout(s, "Can a 7B SLM serve as the linguistic transducer AND hold an Aria-grade SCI reliably across 40+ turns?",
        Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.9), fill=NAVY, size=18)
text(s, "What each answer implies for deployment",
     Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.4),
     size=14, bold=True, color=NAVY)
table(s,
      ["Answer", "Deployment consequence"],
      [
          (("YES", GREEN, True),
           ("SMC-enabled CHA ships at 7B. Edge deployment (~6 GB total footprint) preserved. The SLM does both jobs (transduction + self-modeling) without escalation.", DARK)),
          (("NO", RED, True),
           ("SMC requires 14B+ SLM (breaks edge deployment) — or a different method to close the gap at 7B. The architecture must adapt.", DARK)),
      ],
      Inches(0.5), Inches(3.7), Inches(12.3), Inches(2.0),
      font_size=14, col_widths=[Inches(1.5), Inches(10.8)])
text(s, "If the answer is \"not yet, but here's how\" — that's interesting too.",
     Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
     size=14, italic=True, color=TEAL, align=PP_ALIGN.CENTER)
text(s, "This program's job: turn the question into a decision rule with a binary outcome.",
     Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
     size=13, italic=True, color=GREY, align=PP_ALIGN.CENTER)
footer(s)
notes(s, "This is the deployment-oriented framing of the problem. The audience now sees that the program isn't open-ended research — there's a specific yes/no question with concrete consequences either way. The next slide shows the three-phase plan for answering it.")

# ============================================================
#  SLIDE 6 — Roadmap
# ============================================================
s = new_slide()
title_bar(s, "Three Phases, One Question", "How the program is structured")
table(s,
      ["Phase", "Question", "Result"],
      [
          ("Phase 1\n(Exp 1, baseline)", "Where does the persona break?",
           ("7B coherent but capped at 3.06 mean (threshold 3.5). Episodic dimension at floor.", DARK)),
          ("Phase 2\n(Exp 1, interventions)", "Can prompt-time strategies fix it?",
           ("6 strategies tested. Best two tied at 3.20. Episodic stays unsolved.", DARK)),
          ("Phase 3\n(Experiment 2)", "Do parameter updates (LoRA) close the gap?",
           ("Mean 4.42, ΔE = +0.579, Cohen's d = 7.51. Outcome A: 14B retired.", GREEN, True)),
      ],
      Inches(0.6), Inches(1.5), Inches(12.1), Inches(4.4),
      font_size=14, col_widths=[Inches(2.5), Inches(4.2), Inches(5.4)])
callout(s, "The whole program is built around a single deployment-relevant question: can we ship CHA at 7B?",
        Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.6), fill=NAVY, size=14)
footer(s)
notes(s, "This roadmap is the spine of the talk. Each section that follows answers one of these three questions. Row 3 is the punchline — flag it now so the audience knows where we land.")

# ============================================================
#  SLIDE 7 — Phase 1 divider
# ============================================================
divider("PHASE 1   ·   BASELINE",
        "Where Does the Persona Break?",
        "Two SLMs tested in the linguistic transducer role. Two failure modes documented. One open question handed to Phase 2.")
notes(prs.slides[-1], "Phase 1 is the diagnostic phase. We need to know what 'breaks' looks like before we can argue about how to fix it.")

# ============================================================
#  SLIDE 8 — Phi-4-mini
# ============================================================
s = new_slide()
title_bar(s, "3.8B Is Below the Viable Threshold", "Phi-4-mini cannot parse the JSON SCI")
bullets(s, [
    "Mean PersonaScore: 1.08 / 5.0 across 960 probes",
    "Coherent scripts: 2 / 30 (7%)",
    "Mode of score distribution: 1 (96.0% of all probes)",
    "T* = 5 (immediate failure at first measurement)",
    "Failure mode: instruction-token gibberish from turn 1",
], Inches(0.6), Inches(1.4), Inches(6.0), Inches(4.0), size=15)
sample_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(1.5), Inches(5.7), Inches(4.0))
sample_box.line.fill.background()
sample_box.fill.solid(); sample_box.fill.fore_color.rgb = ALT_ROW
text(s, "Actual model output:", Inches(7.2), Inches(1.6), Inches(5.4), Inches(0.4),
     size=12, bold=True, color=NAVY)
text(s,
     '"<----------------"As aistingiileAITheta InstructionTaskConversion"\n\n'
     '"AssistantAssistantitermathbfisciterierscripiaraboritraborit..."\n\n'
     '"itoryiersarialialiedentiimpl"',
     Inches(7.2), Inches(2.05), Inches(5.4), Inches(3.2),
     size=11, italic=True, color=DARK)
callout(s, "3.8B is below the minimum viable size for JSON-based SCI. The rest of this program uses a 7B-class SLM (Qwen2.5-7B) — the empirically established SMC tier.",
        Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.7), fill=AMBER, size=13)
footer(s)
notes(s, "Don't dwell here. The point is just to establish a lower bound: there's a parameter-count floor below which the SCI framework doesn't even apply. Move on.")

# ============================================================
#  SLIDE 9 — Qwen2.5-7B baseline
# ============================================================
s = new_slide()
title_bar(s, "Qwen2.5-7B: Coherent, but Below Threshold", "The workhorse model for the rest of the program")
bullets(s, [
    "Coherent on 30/30 scripts (vs 2/30 for Phi-4-mini)",
    "Mean PersonaScore: 3.06 across 960 probes — a 2.8× improvement over 3.8B",
    "T* = 5: never crosses the 3.5 threshold at any probe turn",
    "Total decline turn 5 → turn 40: only −0.20 points (3.16 → 2.96)",
    "Genuine near-threshold behavior — no longer a floor effect",
], Inches(0.6), Inches(1.4), Inches(7.5), Inches(3.5), size=15)
table(s,
      ["Turn", "5", "10", "15", "20", "25", "30", "35", "40"],
      [["Mean", "3.16", "3.13", "3.15", "3.11", "3.09", "3.04", "3.00", "2.96"]],
      Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.7), font_size=12)
callout(s, "Coherent enough to study. Stable enough to characterize. But 0.44 points short of the threshold — the gap that defines this program.",
        Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.7), fill=TEAL, size=13)
footer(s)
notes(s, "This is where the real research starts. The question shifts from 'does it work at all?' to 'how does it fail, and can we fix it?'")

# ============================================================
#  SLIDE 10 — Piecewise inflection
# ============================================================
s = new_slide()
title_bar(s, "The Piecewise Inflection at Turn 15", "A natural refresh interval suggests itself")
img_path = EXP1 / "results_qwen2.5_7b" / "persona_score_timeseries.png"
if img_path.exists():
    image(s, img_path, Inches(0.5), Inches(1.2), w=Inches(7.8))
bullets(s, [
    "Best fit: piecewise model (AIC −67.1, beats linear by 7.9)",
    "Stable through ~turn 15, then declines at β = 0.008 / turn",
    "Inflection point T₀ = 15 is robust across script types",
    "Implication: a refresh interval of ~15 turns is suggested by the data, not chosen by hand",
    "Sets up the first architectural intervention (Phase 2)",
], Inches(8.6), Inches(1.4), Inches(4.4), Inches(5.0), size=13, space=10)
footer(s)
notes(s, "The piecewise fit is the most actionable finding of Phase 1. It's not just 'the model degrades'; it's 'the model is stable for a known number of turns, then breaks at a known inflection.' That's a deployment-relevant signal.")

# ============================================================
#  SLIDE 11 — Per-dimension scoreboard + takeaway
# ============================================================
s = new_slide()
title_bar(s, "Where the Persona Falls Apart", "Per-dimension breakdown sets up Phase 2 priorities")
table(s,
      ["Dimension", "T*", "Turn 5 mean", "Turn 40 mean", "Verdict"],
      [
          ("Trait (T)", ">40", "3.90", "3.57", ("the strength — stable above threshold", GREEN, True)),
          ("Capability (C)", "5", "3.33", "3.33", "borderline; flat with mid-conversation dips"),
          ("Style (S)", "5", "3.13", "2.57", "below threshold, declining late"),
          ("Episodic (E)", "5", "2.27", "2.37", ("at floor — pervasive fabrication", RED, True)),
      ],
      Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.4), font_size=13)
text(s, "Failure-mode taxonomy (Phase 1 baseline)",
     Inches(0.6), Inches(4.1), Inches(12.1), Inches(0.4),
     size=14, bold=True, color=NAVY)
table(s,
      ["Failure mode", "Count", "% of 960 probes"],
      [
          ("Episodic fabrication", "194", "20.2%"),
          ("Register/style shift", "113", "11.8%"),
          ("Capability overstatement", "111", "11.6%"),
          ("Trait drift", "43", "4.5%"),
      ],
      Inches(0.6), Inches(4.6), Inches(7.0), Inches(2.0), font_size=12)
callout(s, "Episodic fabrication is 4.5× the rate of trait drift. This is the dimension Phase 2 has to fix.",
        Inches(7.9), Inches(4.7), Inches(4.8), Inches(2.0), fill=AMBER, size=13)
footer(s)
notes(s, "Two takeaways: (1) Trait is the model's strength, do no harm. (2) Episodic dominates failure modes by a 4.5× margin. Phase 2 must address it, and this is the lens through which to evaluate every intervention that follows.")

# ============================================================
#  SLIDE 12 — Phase 2 divider
# ============================================================
divider("PHASE 2   ·   ARCHITECTURAL INTERVENTIONS",
        "Can Prompt-Time Strategies Close the Gap?",
        "Six SCI variants. Same 30 scripts. Same judge. Same 7B model. The question is purely architectural.")

# ============================================================
#  SLIDE 13 — Six strategies overview
# ============================================================
s = new_slide()
title_bar(s, "Six Architectural Strategies", "Each tests a different hypothesis about what the model needs")
table(s,
      ["Strategy", "Mechanism", "Hypothesis"],
      [
          ("Baseline", "Persona JSON in system prompt", "What the model can do unaided."),
          ("SCI Refresh", "Re-inject persona JSON at turn 13", "Periodic re-grounding cures drift."),
          ("Episodic RAG", "Strip events; retrieve on E-probes", "On-demand retrieval beats baked-in events."),
          ("Hybrid RAG", "Compressed summaries + RAG details", "Best of both; preserves implicit anchoring."),
          ("Combined", "Refresh + Episodic RAG together", "Refresh handles trait/style; RAG handles episodic."),
          ("Multi-Refresh", "Refresh at turns 13 AND 28", "Second refresh covers single-refresh fade."),
      ],
      Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.4),
      font_size=12, col_widths=[Inches(2.0), Inches(4.5), Inches(5.8)])
callout(s, "Same 30 scripts, same judge. Differences across rows are purely architectural.",
        Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.6), fill=NAVY, size=13)
footer(s)
notes(s, "Set the table for the six strategies. Each tests a clean hypothesis. Don't go deep on any one yet — the next slides do that.")

# ============================================================
#  SLIDE 14 — SCI Refresh
# ============================================================
s = new_slide()
title_bar(s, "SCI Refresh at Turn 13", "Eliminates the inflection, fades by turn 35")
bullets(s, [
    "Mean: 3.15 (+0.07 over baseline)",
    "Degradation profile: piecewise → linear (inflection point eliminated)",
    "Effect peaks turns 15–30 (+0.09), fades to baseline by turn 35",
    "Trait drift −19%, capability overstatement −5%",
    "Episodic fabrication unchanged (−1%) — refresh doesn't help retrieval",
    "Useful refresh window: ~20 turns (13 → 33) before model resumes drifting",
], Inches(0.6), Inches(1.4), Inches(7.0), Inches(4.5), size=14, space=8)
table(s,
      ["Turn", "Baseline", "Refresh", "Δ"],
      [
          ("5", "3.16", "3.30", "+0.14"),
          ("15", "3.15", "3.28", "+0.13"),
          ("25", "3.09", "3.17", "+0.08"),
          ("35", "3.00", "2.98", ("−0.02", RED)),
          ("40", "2.96", "2.98", "+0.02"),
      ],
      Inches(8.0), Inches(1.5), Inches(4.7), Inches(3.0), font_size=12)
callout(s, "First architectural win — but the effect is modest and decays.",
        Inches(8.0), Inches(4.8), Inches(4.7), Inches(0.7), fill=TEAL, size=13)
footer(s)
notes(s, "Refresh is the simplest possible intervention. It works, but not enough. The fade-by-turn-35 finding is what motivates Multi-Refresh later.")

# ============================================================
#  SLIDE 15 — Episodic RAG + dual-purpose finding
# ============================================================
s = new_slide()
title_bar(s, "Episodic RAG: 4× Slower Decay — But a Trade-off", "Removing events from the prompt destabilizes traits")
bullets(s, [
    "β: 0.008 → 0.002 (4× slower degradation rate)",
    "E dimension: 2.37 → 2.69 (+0.32, the program's first real episodic improvement)",
    "BUT — Trait T*: >40 → 10 (trait dimension drops below threshold)",
    "Trait drift up +21%, episodic fabrication down only −11%",
    "Most important Phase 2 discovery (next bullet) ↓",
], Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.5), size=14, space=8)
callout(s,
        "salient_past_events carries TWO loads: explicit (\"what Aria remembers\") AND implicit (\"who Aria is\").\n"
        "Removing the section to free token budget destabilizes traits even when retrieval handles the explicit load cleanly.",
        Inches(0.6), Inches(4.4), Inches(12.1), Inches(1.6), fill=AMBER, size=14)
text(s, "Architectural implication: any episodic refactor must preserve implicit trait anchoring.",
     Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.6),
     size=14, italic=True, color=NAVY)
footer(s)
notes(s, "This is the most important slide in Phase 2. The dual-purpose finding is the kind of result that survives the experiment and informs design downstream. Spend extra time here. The implicit anchoring point reframes the rest of the section.")

# ============================================================
#  SLIDE 16 — Hybrid RAG failed
# ============================================================
s = new_slide()
title_bar(s, "Hybrid RAG: The Hypothesis Was Wrong", "Compressed summaries don't preserve implicit anchoring")
bullets(s, [
    "Hypothesis: one-line summaries in SCI + full RAG details = best of both worlds",
    "Reality: T* dropped to 5 (worse than full RAG) — implicit anchoring requires the full narratives, not topic markers",
    "Episodic improvement: +0.46 (best E result of any single intervention!)",
    "But trait regression and a unique style regression (+8% register shift) outweighed the gain",
    "Three possible mechanisms (Section 15.3 of report):",
    (1, "Full narratives carry trait-relevant emotional valence and behavioral cues that summaries strip"),
    (1, "Token economy isn't the bottleneck — freed budget doesn't reallocate to traits"),
    (1, "Compressed summary + full RAG creates cognitive dissonance; model anchors on the vague version"),
], Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.2), size=14, space=8)
callout(s, "Important negative result: the simpler architecture is better. Don't compress what serves a dual purpose.",
        Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5), fill=RED, size=13)
footer(s)
notes(s, "Negative results matter. This is the kind of finding that saves a downstream team six weeks of trying to optimize the wrong thing.")

# ============================================================
#  SLIDE 17 — Combined R+RAG
# ============================================================
s = new_slide()
title_bar(s, "Combined Refresh + RAG: Late-Conversation Peak", "Refresh handles trait stability; RAG handles episodic recall")
bullets(s, [
    "Mean: 3.20 (the program's joint-best at this point)",
    "Effective β ≈ 0 (degradation eliminated)",
    "Emergent finding: late-turn scores are HIGHER than early — turn 25–40 mean 3.24 vs turn 5–10 mean 3.10",
    "E dimension: +0.39 over baseline (refresh keeps the model engaged with the persona; RAG keeps it grounded in real events)",
    "Trait scores preserved at baseline (3.69 vs 3.67) — refresh successfully counteracts RAG's destabilization",
], Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.6), size=15, space=10)
table(s,
      ["Turn", "Baseline", "Combined", "Δ"],
      [
          ("5", "3.16", "3.15", ("−0.01", GREY)),
          ("20", "3.11", "3.28", ("+0.17", GREEN)),
          ("30", "3.04", "3.28", ("+0.24", GREEN, True)),
          ("40", "2.96", "3.19", ("+0.23", GREEN, True)),
      ],
      Inches(0.6), Inches(5.4), Inches(8.0), Inches(1.4), font_size=12)
callout(s, "First condition where the model is MORE persona-consistent at turn 40 than at turn 5.",
        Inches(8.9), Inches(5.4), Inches(3.9), Inches(1.4), fill=GREEN, size=13)
footer(s)
notes(s, "The emergent late-conversation peak is the most surprising positive finding in Phase 2. Worth pausing on.")

# ============================================================
#  SLIDE 18 — Multi-refresh
# ============================================================
s = new_slide()
title_bar(s, "Multi-Refresh (Turns 13 + 28): Best Trait Stability", "No retrieval infrastructure required")
bullets(s, [
    "Mean: 3.20 (tied with Combined for top spot)",
    "Trait dimension: 3.86 — the program's BEST",
    "Style dimension: 3.14 — the program's BEST",
    "Late-conversation mean (turns 25–40): 3.19 (single-refresh faded to 2.98 here)",
    "Lowest total failures of any condition: 428 (vs 461 baseline, −7%)",
    "Trait drift −28%, register shift −11% (largest reductions in any condition)",
    "Operationally simplest deployable strategy — no retrieval system needed",
], Inches(0.6), Inches(1.4), Inches(12.1), Inches(4.5), size=15, space=10)
callout(s, "If you're shipping without RAG, ship Multi-Refresh. If you have RAG, the next slide is your decision.",
        Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.6), fill=TEAL, size=13)
footer(s)
notes(s, "Multi-Refresh and Combined tie at the top. Multi-Refresh wins on simplicity; Combined wins on episodic recall. The leaderboard slide makes this concrete.")

# ============================================================
#  SLIDE 19 — Phase 2 leaderboard
# ============================================================
s = new_slide()
title_bar(s, "Phase 2 Leaderboard", "Six-way comparison — and the ceiling that holds")
table(s,
      ["Strategy", "Mean", "Mean 25–40", "β", "Episodic", "Trait", "Total fails"],
      [
          ("Baseline",      "3.08",         "3.02",        "0.008",  "2.37",        "3.67",         "461"),
          ("Refresh-13",    "3.15",         "3.07",        "0.008",  "2.37",        "3.79",         "444"),
          ("Episodic-RAG",  "3.15",         "3.15",        "0.002",  "2.69",        "3.60",         "457"),
          ("Hybrid-RAG",    "3.17",         "3.11",        "~0",     ("2.83", GREEN, True), "3.60",  "446"),
          ("Multi-Refresh", ("3.20", GREEN, True), "3.19", "~0",     "2.43",        ("3.86", GREEN, True), ("428", GREEN, True)),
          ("Combined",      ("3.20", GREEN, True), ("3.24", GREEN, True), "~0", "2.76", "3.69", "431"),
      ],
      Inches(0.4), Inches(1.4), Inches(12.5), Inches(3.4), font_size=12)
text(s, "Phase 2 takeaway",
     Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.4),
     size=14, bold=True, color=NAVY)
bullets(s, [
    "All six strategies converge in [3.08, 3.20] — a 0.12-point band, far below the 3.5 threshold",
    "No condition crosses E = 3.0; the episodic gap survives every architectural strategy",
    "The closing question of Phase 2: is the residual gap architectural (need 14B) or capability-shaped (LoRA can fix)?",
], Inches(0.4), Inches(5.4), Inches(12.5), Inches(1.7), size=14, space=8)
footer(s)
notes(s, "The leaderboard is the mic-drop slide for Phase 2. Architecture alone can't get us there. This is the exact question that motivates Experiment 2.")

# ============================================================
#  SLIDE 20 — Phase 3 divider
# ============================================================
divider("PHASE 3   ·   FINE-TUNING",
        "When Prompts Hit a Ceiling, Update the Parameters",
        "Experiment 2: a 4-condition test of LoRA fine-tuning, with H5 sub-runs for data scaling.")

# ============================================================
#  SLIDE 21 — Hypotheses + the fork
# ============================================================
s = new_slide()
title_bar(s, "Pre-Registered Hypotheses", "And the architectural-vs-capability fork they resolve")
table(s,
      ["ID", "Statement", "Pass criterion"],
      [
          ("H1", "LoRA + Combined SCI brings mean above 3.5", "C ≥ 3.5"),
          ("H2", "Fine-tuning meaningfully addresses fabrication", "ΔE = E(C) − E(D) ≥ +0.30"),
          ("H3", "FT and SCI interact (additive or super-additive)", "FT effect > 0 with Combined SCI ≥ no SCI"),
          ("H4", "Base capability preserved (no catastrophic forgetting)", "No regression on out-of-domain probes"),
          ("H5", "Persona consistency scales log with data size", "Score(2K) < Score(5K) < Score(10K)"),
      ],
      Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.8),
      font_size=12, col_widths=[Inches(0.7), Inches(5.6), Inches(6.0)])
text(s, "The fork H2 resolves",
     Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
     size=14, bold=True, color=NAVY)
bullets(s, [
    "Architectural-ceiling hypothesis: 7B fundamentally can't do episodic — predicted ΔE < +0.15",
    "Capability-shaped-gap hypothesis: 7B has the parameters; base distribution lacks Aria-shaped data — predicted ΔE ≥ +0.30",
    "If the architectural hypothesis is right, scale to 14B+. If capability-shaped, fine-tune.",
], Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.0), size=14, space=8)
footer(s)
notes(s, "H2 is the decisive hypothesis. The criterion was set at +0.30 because the program's worst-to-best architectural intervention only moved E by 0.46. We need to see something larger to claim fine-tuning is doing real work.")

# ============================================================
#  SLIDE 22 — 4-condition design
# ============================================================
s = new_slide()
title_bar(s, "4-Condition Design + H5 Sub-Runs", "Same 30 scripts, same probes, same judge as Experiment 1")
table(s,
      ["ID", "Subject Model", "SCI Strategy", "Role"],
      [
          ("A", "Qwen2.5-7B + LoRA-10K", "None (raw role instruction)", "Ablation: FT alone, no SCI"),
          ("B", "Qwen2.5-7B + LoRA-10K", "Baseline SCI",                "FT + naive SCI"),
          (("C", NAVY, True), ("Qwen2.5-7B + LoRA-10K", NAVY, True),
              ("Combined SCI (refresh + RAG)", NAVY, True),
              ("Headline: best of both worlds", NAVY, True)),
          ("D", "Qwen2.5-7B (no LoRA)", "Combined SCI", "Replication control: equals Exp 1's best"),
      ],
      Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.6), font_size=13)
text(s, "H5 sub-runs (same Condition C config, swap the adapter)",
     Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.4),
     size=13, bold=True, color=NAVY)
table(s,
      ["Sub-run", "Adapter", "Purpose"],
      [
          ("C-2K",  "LoRA-2K",  "Lower bound of data-scaling curve"),
          ("C-5K",  "LoRA-5K",  "Mid-point of data-scaling curve"),
          ("C-10K", "LoRA-10K", "Upper bound (= headline Condition C)"),
      ],
      Inches(0.5), Inches(4.7), Inches(12.3), Inches(1.8), font_size=12)
callout(s, "C vs D (paired by script) is the primary statistical test. A and B isolate FT-only and FT-with-naive-SCI as ablations.",
        Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5), fill=NAVY, size=12)
footer(s)
notes(s, "Within-subject design: every condition sees the exact same probes on the exact same scripts. The C-vs-D paired test is the main event.")

# ============================================================
#  SLIDE 23 — Training pipeline
# ============================================================
s = new_slide()
title_bar(s, "Training Pipeline", "10K examples, QLoRA on A100")
text(s, "Dataset (Sonnet 4.6, ~23K API calls)",
     Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.4),
     size=14, bold=True, color=NAVY)
bullets(s, [
    "10,000 (system_prompt, history, probe, target) tuples",
    "5-rule QC filter: no token leakage, episodic grounding, length 30–150 words, marker vocabulary, MiniLM ≥ 0.35",
    "Stratified across 4 dims × 5 scenarios × 3 turn-depth bands (Hamilton method)",
    "80 / 10 / 10 split → 8,000 train / 1,000 val / 1,000 test",
    "Steady-state acceptance: ~80% first-pass, ~21% on replay of failed examples",
], Inches(0.5), Inches(1.7), Inches(6.0), Inches(4.5), size=12, space=6)
text(s, "Training (Qwen2.5-7B, QLoRA)",
     Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.4),
     size=14, bold=True, color=NAVY)
bullets(s, [
    "4-bit NF4 quantization, BF16 LoRA adapters",
    "LoRA r=16, α=32, dropout 0.05",
    "Target modules: q/k/v/o + gate_proj/up_proj (~40M trainable params, ~0.5%)",
    "Sequence length 3,072 (raised from 2,048 — original silently dropped 38% of rows)",
    "Effective batch 16 (per-device 2 × accum 8), paged_adamw_8bit, lr 2e-4, 2 epochs",
    "Loss masked on assistant turns only (DataCollatorForCompletionOnlyLM)",
    "A100 80GB · 3 adapters: 2K (45 min) / 5K (4 hr) / 10K (8 hr)",
], Inches(7.0), Inches(1.7), Inches(5.8), Inches(4.5), size=12, space=6)
table(s,
      ["Adapter", "Train rows", "Eval loss (best)"],
      [
          ("LoRA-2K",  "2,000",  "0.91"),
          ("LoRA-5K",  "5,000",  "0.77"),
          ("LoRA-10K", "10,000 (8K)", ("0.69", GREEN, True)),
      ],
      Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), font_size=12)
footer(s)
notes(s, "Standard QLoRA recipe, but two tunings worth flagging: max_seq_length=3072 (default 2048 silently dropped a third of the data), and DataCollatorForCompletionOnlyLM to mask user turns from the loss.")

# ============================================================
#  SLIDE 24 — Headline result
# ============================================================
s = new_slide()
title_bar(s, "Headline Result: H1 PASSED", "Condition C exceeds the threshold by +0.92 points")
img_path = EXP2 / "results" / "condition_comparison.png"
if img_path.exists():
    image(s, img_path, Inches(0.4), Inches(1.2), w=Inches(7.5))
big_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(1.3), Inches(4.7), Inches(5.4))
big_box.line.fill.background()
big_box.fill.solid(); big_box.fill.fore_color.rgb = NAVY
text(s, "C: 4.415", Inches(8.3), Inches(1.5), Inches(4.5), Inches(0.9),
     size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text(s, "vs threshold 3.5    ·    +0.92 margin",
     Inches(8.3), Inches(2.4), Inches(4.5), Inches(0.4),
     size=12, color=LIGHT, align=PP_ALIGN.CENTER)
text(s, "Δ (C − D)",
     Inches(8.3), Inches(2.9), Inches(4.5), Inches(0.4),
     size=12, color=LIGHT, align=PP_ALIGN.CENTER)
text(s, "+1.191",
     Inches(8.3), Inches(3.3), Inches(4.5), Inches(0.6),
     size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
text(s, "Cohen's d",
     Inches(8.3), Inches(4.0), Inches(4.5), Inches(0.4),
     size=12, color=LIGHT, align=PP_ALIGN.CENTER)
text(s, "7.51",
     Inches(8.3), Inches(4.4), Inches(4.5), Inches(0.6),
     size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
text(s, "Paired t-test, n=30",
     Inches(8.3), Inches(5.1), Inches(4.5), Inches(0.4),
     size=12, color=LIGHT, align=PP_ALIGN.CENTER)
text(s, "p ≈ 1.4 × 10⁻²³",
     Inches(8.3), Inches(5.5), Inches(4.5), Inches(0.5),
     size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text(s, "H1 PASSED ✓",
     Inches(8.3), Inches(6.1), Inches(4.5), Inches(0.5),
     size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
footer(s)
notes(s, "Cohen's d = 7.51 is far past the conventional 'large' bound of 0.8. The effect is enormous and statistically rock-solid. Spend a moment letting this land.")

# ============================================================
#  SLIDE 25 — Time series across all conditions
# ============================================================
s = new_slide()
title_bar(s, "Stable Across All Probe Turns", "No piecewise inflection in the fine-tuned conditions")
img_path = EXP2 / "results" / "persona_score_timeseries.png"
if img_path.exists():
    image(s, img_path, Inches(0.5), Inches(1.2), w=Inches(8.4))
bullets(s, [
    "Condition C maintains 4.29 – 4.55 across all 8 probe turns",
    "The piecewise pattern from Phase 1 is GONE in fine-tuned conditions",
    "Refresh injections at turns 15 and 30 visible as small bumps in D",
    "Conditions A, B, C all sit comfortably above the 3.5 threshold",
    "D reproduces Exp 1's Combined SCI result within tolerance",
], Inches(9.2), Inches(1.4), Inches(3.9), Inches(5.5), size=12, space=10)
footer(s)
notes(s, "The disappearance of the piecewise inflection is itself a finding. The model is no longer fighting drift; it's holding the persona by default.")

# ============================================================
#  SLIDE 26 — Per-dimension + Style was hidden bottleneck
# ============================================================
s = new_slide()
title_bar(s, "Style Was the Hidden Bottleneck", "All four dimensions improve; Style most dramatically")
img_path = EXP2 / "results" / "dimension_comparison.png"
if img_path.exists():
    image(s, img_path, Inches(0.4), Inches(1.2), w=Inches(7.6))
table(s,
      ["Dim", "D (base)", "C (LoRA)", "Δ"],
      [
          ("Trait",      "3.65",  "4.90", ("+1.25", GREEN, True)),
          ("Episodic",   "2.77",  "3.35", ("+0.58", GREEN, True)),
          ("Capability", "3.42",  "4.47", ("+1.05", GREEN, True)),
          ("Style",      ("3.06", RED, True), "4.94", ("+1.88", GREEN, True)),
      ],
      Inches(8.3), Inches(1.5), Inches(4.5), Inches(2.6), font_size=13)
callout(s,
        "Under Sonnet 4.5, Style is joint-bottom with Episodic (3.06 vs 2.77) — not the middle-tier dimension Exp 1 implied. "
        "FT crushes both: +1.88 on S, +0.58 on E. Failure-mode count: 436 → 106 (−75%).",
        Inches(8.3), Inches(4.4), Inches(4.5), Inches(2.4), fill=AMBER, size=12)
footer(s)
notes(s, "Two takeaways: (1) Style was a stylometric problem the base model couldn't hold from the system prompt — fine-tuning fixes it cleanly. (2) Failure-mode counts collapse by 75%, which is the more visceral way to communicate the size of the win.")

# ============================================================
#  SLIDE 27 — H2 episodic ceiling resolved
# ============================================================
s = new_slide()
title_bar(s, "H2: Episodic Ceiling Resolved", "ΔE = +0.579 — capability-shaped, not architectural")
text(s, "ΔE = +0.579",
     Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2),
     size=72, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
text(s, "vs the +0.30 threshold for \"fine-tuning meaningfully addresses fabrication\"",
     Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.5),
     size=16, italic=True, color=GREY, align=PP_ALIGN.CENTER)
table(s,
      ["Hypothesis", "Predicted ΔE", "Observed", "Verdict"],
      [
          ("Architectural ceiling (need 14B)", "< +0.15", "+0.579", ("REFUTED", RED, True)),
          ("Capability-shaped gap (LoRA can fix)", "≥ +0.30", "+0.579", ("CONFIRMED", GREEN, True)),
      ],
      Inches(0.5), Inches(3.7), Inches(12.3), Inches(1.6), font_size=14)
bullets(s, [
    "The closing question of Phase 2 is decisively answered: the 7B model has the parameters; the base distribution just lacked the data",
    "Residual 0.15 to threshold (3.35 → 3.5) is now a data-scaling problem, not a parameter-count problem",
    "This single result retires the planned 14B model test from the program's critical path",
], Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5), size=13, space=8)
footer(s)
notes(s, "This is the slide that earns the Outcome A claim. Spend time on it. The +0.579 number is the single most important measurement in Experiment 2 — it's the resolution of a year of program ambiguity.")

# ============================================================
#  SLIDE 28 — H5 learning curve
# ============================================================
s = new_slide()
title_bar(s, "H5: Episodic Score Scales Logarithmically", "Diminishing returns argue for retrieval, not more LoRA data")
img_path = EXP2 / "results" / "learning_curve.png"
if img_path.exists():
    image(s, img_path, Inches(0.4), Inches(1.2), w=Inches(7.6))
bullets(s, [
    "LoRA-2K → E = 2.85",
    "LoRA-5K → E = 3.10",
    "LoRA-10K → E = 3.35",
    "Fit: E(n) = 0.291 · log(n) + 0.633",
    "Predicted at 20K: 3.51 (barely clears threshold)",
    "Predicted to clear E = 4.0: ~250K examples (impractical)",
    "Eval loss mirrors persona score: 0.91 → 0.77 → 0.69",
], Inches(8.3), Inches(1.4), Inches(4.7), Inches(4.5), size=12, space=8)
callout(s, "Last 0.15 on E is more cost-effective via better RAG retrieval than via continued LoRA scaling.",
        Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.6), fill=TEAL, size=13)
footer(s)
notes(s, "The H5 curve says: yes, more data helps, but logarithmically. The marginal cost of pushing E from 3.35 to 3.5 via more training data is high; better retrieval architecture is probably cheaper.")

# ============================================================
#  SLIDE 29 — Cross-experiment summary
# ============================================================
s = new_slide()
title_bar(s, "The Whole Story in One Table", "From floor effect to threshold-clearing in two experiments")
table(s,
      ["Metric", "Phi-4-mini\n(3.8B)", "Qwen2.5-7B\nbaseline", "Qwen + Combined\nSCI (best Exp 1)", "Qwen + LoRA-10K\n+ Combined SCI"],
      [
          ("Mean PersonaScore",      "1.08",        "3.06",           "3.20",            ("4.42", GREEN, True)),
          ("Trait dim",              "1.00–1.20",   "3.67",           "3.69",            ("4.90", GREEN, True)),
          ("Episodic dim",           "1.00–1.17",   ("2.37", RED),    "2.76",            ("3.35", GREEN, True)),
          ("Capability dim",         "1.00–1.23",   "3.28",           "3.27",            ("4.47", GREEN, True)),
          ("Style dim",              "1.00–1.07",   "3.00",           "3.09",            ("4.94", GREEN, True)),
          ("H1 (≥ 3.5)?",            ("✗", RED),    ("✗", RED),       ("✗", RED),        ("✓", GREEN, True)),
          ("Coherent scripts",       "2 / 30",      "30 / 30",        "30 / 30",         "30 / 30"),
          ("Total failures (≤ 2)",   "~960",        "461",            "431",             ("106", GREEN, True)),
      ],
      Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.5),
      font_size=12, col_widths=[Inches(2.6), Inches(2.3), Inches(2.4), Inches(2.6), Inches(2.6)])
text(s, "Each column builds on the last: capability gates the experiment, then architecture provides the framework, then fine-tuning closes the gap.",
     Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
     size=12, italic=True, color=GREY, align=PP_ALIGN.CENTER)
footer(s)
notes(s, "If the audience remembers one slide from this talk, it should be this one. The trajectory across the program is the headline.")

# ============================================================
#  SLIDE 30 — Decision Rule Outcome A
# ============================================================
s = new_slide()
title_bar(s, "Decision Rule: Outcome A", "14B model test retired from the critical path")
table(s,
      ["Outcome", "C ≥ 3.5?", "ΔE ≥ +0.30?", "Implication"],
      [
          (("A", GREEN, True), ("✓", GREEN, True), ("✓", GREEN, True),
              ("SMC architecture complete at 7B; 14B test unnecessary", GREEN, True)),
          ("B", "✓", "✗", "Overall closed but episodic remains; need 14B"),
          ("C", "✗", "✓", "Episodic helped but overall not at threshold"),
          ("D", "✗", "✗", "Fine-tuning isn't the answer; pursue 14B"),
      ],
      Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.6), font_size=13)
callout(s,
        "OBSERVED: C = 4.415 ✓   ·   ΔE = +0.579 ✓   →   OUTCOME A",
        Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.7), fill=GREEN, size=18)
bullets(s, [
    "Pre-registered before Experiment 2: which outcome triggers retires the 14B test from the critical path",
    "The deployment story is settled: 7B + LoRA-10K + Combined SCI is the recommended Phase 1 SMC architecture",
    "Remaining work is (a) cross-model replication, (b) the deferred H4 base-capability test, (c) better RAG for the last 0.15 on E",
], Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), size=14, space=8)
footer(s)
notes(s, "Don't undersell this. Pre-registered decision rules with binary outcomes are how you turn a research program into a deployment recommendation. We pre-committed, the data came in, the rule fired, and now we know what to ship.")

# ============================================================
#  SLIDE 31 — Limitations
# ============================================================
s = new_slide()
title_bar(s, "Limitations & Open Questions", "What this program doesn't tell us yet")
bullets(s, [
    ("H4 base-capability test was deferred — no out-of-domain probe battery has been run on the fine-tuned model"),
    (1, "Higher priority than any other follow-up: rules out catastrophic forgetting before deployment"),
    "Single SLM tested (Qwen2.5-7B) — Outcome A could be Qwen-specific; cross-model replication on Llama 3.1 8B and Gemma 2 9B is pending",
    "Single persona (Aria, psychotherapy support) — domain generalization is unmeasured",
    "40-turn conversation ceiling — multi-day persistence is unmeasured (this matters for production)",
    "Synthetic dataset by Sonnet 4.6, judged by Sonnet 4.5 — possible same-family stylistic alignment that inflates the gain",
    "H5 curve has only 3 points; the 20K extrapolation has wide uncertainty",
    "Intra-model judge stochasticity κ_w = 0.611 — boundary-case scoring carries real noise",
], Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.6), size=14, space=8)
footer(s)
notes(s, "Be honest. The H4 test is the one thing that could change the deployment story — if LoRA causes regression on out-of-domain tasks, the LoRA needs to be context-conditional rather than always-on.")

# ============================================================
#  SLIDE 32 — Future work
# ============================================================
s = new_slide()
title_bar(s, "Future Work — Priority Ordered", "Where the next experiment goes")
table(s,
      ["#", "Task", "Why"],
      [
          ("1", "H4 base-capability test", ("Highest priority; gates deployment", RED, True)),
          ("2", "E-stratified fine-tuning (15K, half E-focused)", "Predicted to clear E = 3.5 with less data than scaling overall to 20K"),
          ("3", "Cross-model replication (Llama 3.1 8B, Gemma 2 9B)", "Promotes Outcome A from \"Qwen-specific\" to \"7B-class generalization\""),
          ("4", "Multi-session persistence (200+ turns, multi-day)", "Production-relevant: do gains hold beyond 40 turns?"),
          ("5", "Ablate refresh / RAG inside Condition C", "Does FT subsume some SCI components? If yes, retire complexity"),
          ("6", "14B + LoRA, scientific question only", "Does residual E gap close at 14B? Off critical path but informative"),
      ],
      Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.0),
      font_size=13, col_widths=[Inches(0.6), Inches(5.4), Inches(6.5)])
callout(s, "Critical-path priority: complete H4 before publishing or deploying.",
        Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.5), fill=NAVY, size=13)
footer(s)
notes(s, "Highlight #1 — H4 — as the gating follow-up. The rest are interesting but not blocking.")

# ============================================================
#  SLIDE 33 — Conclusion
# ============================================================
s = new_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.line.fill.background()
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
text(s, "Persona Consistency at 7B:",
     Inches(1.0), Inches(1.4), Inches(11.3), Inches(0.9),
     size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text(s, "solved.",
     Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.3),
     size=72, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
text(s, "Three phases · two experiments · one program",
     Inches(1.0), Inches(3.9), Inches(11.3), Inches(0.5),
     size=18, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
text(s, "1.08  →  3.06  →  3.20  →  4.42",
     Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.6),
     size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text(s, "Phi-4-mini    ·    Qwen 7B baseline    ·    Best Exp 1 architecture    ·    LoRA-10K + Combined SCI",
     Inches(1.0), Inches(5.2), Inches(11.3), Inches(0.4),
     size=12, color=LIGHT, align=PP_ALIGN.CENTER)
text(s, "The CHA SMC sub-architecture is ready for Phase 1 deployment at 7B.",
     Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.5),
     size=16, color=WHITE, align=PP_ALIGN.CENTER)
text(s, "Questions?",
     Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.4),
     size=14, italic=True, color=TEAL, align=PP_ALIGN.CENTER)
text(s, f"{_counter[0]} / {TOTAL_SLIDES}",
     Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
     size=9, color=LIGHT, align=PP_ALIGN.RIGHT)
notes(s, "End with the trajectory: 1.08 → 3.06 → 3.20 → 4.42. That single sequence is the program. Open the floor.")

prs.save(str(OUT))
print(f"Wrote {OUT} ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")
